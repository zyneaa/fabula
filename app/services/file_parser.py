from pathlib import Path

import fitz
from docx import Document
from pptx import Presentation


def parse_file(file_path: str) -> str:
    path = Path(file_path)
    ext = path.suffix.lower()

    if ext == ".pdf":
        return _parse_pdf(file_path)
    elif ext == ".docx":
        return _parse_docx(file_path)
    elif ext == ".pptx":
        return _parse_pptx(file_path)
    elif ext == ".txt":
        return path.read_text(encoding="utf-8")
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def _parse_pdf(file_path: str) -> str:
    doc = fitz.open(file_path)
    text = "".join(str(page.get_text()) for page in doc)
    doc.close()
    return text


def _parse_docx(file_path: str) -> str:
    doc = Document(file_path)
    return "\n".join(para.text for para in doc.paragraphs)


def _parse_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(getattr(shape, "text"))
    return "\n".join(text)
